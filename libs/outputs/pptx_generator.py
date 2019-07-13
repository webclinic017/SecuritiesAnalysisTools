from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

import pandas as pd 
import numpy as np 
from datetime import datetime
import os 
import glob 

from libs.utils import fund_list_extractor, windows_compatible_file_parse

# Slide Layouts
PRES_TITLE_SLIDE = 0
TITLE_CONTENT_SLIDE = 1
SECTION_HEADER_SLIDE = 2
TWO_CONTENT_SLIDE = 3
COMPARISON_SLIDE = 4
TITLE_ONLY_SLIDE = 5
BLANK_SLIDE = 6
CONTENT_W_CAPTION_SLIDE = 7
PICTURE_W_CAPTION_SLIDE = 8


def title_presentation(year: str, VERSION: str, wide_ratio=True):
    prs = Presentation()

    height = prs.slide_height
    width = int(16 * height / 9)
    prs.slide_width = width
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
    # else:
    #     slide = prs.slides.add_slide(prs.slide_layouts[PRES_TITLE_SLIDE])

    # title = slide.shapes.title
    LEFT_INCHES = 6
    left = Inches(LEFT_INCHES)
    top = Inches(2.45)
    text = slide.shapes.add_textbox(left, top, Inches(1), Inches(1))
    text_frame = text.text_frame

    # text_frame = title.text_frame
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = f'Securities Analysis'
    p.font.bold = True
    p.font.size = Pt(48)
    p.font.name = 'Arial'

    p4 = text_frame.add_paragraph()
    p4.alignment = PP_ALIGN.CENTER
    p4.text = f"A Technical Analysis of Financial Markets by 'nga-27'"
    p4.font.italic = True
    p4.font.size = Pt(14)
    p4.font.color.rgb = RGBColor(0x74, 0x3c, 0xe6)
    p4.font.name = 'Arial'

    left = Inches(LEFT_INCHES)
    top = Inches(4.0)
    text = slide.shapes.add_textbox(left, top, Inches(1), Inches(1))
    text_frame2 = text.text_frame
    # else:
    #     stitle = slide.placeholders[1]
    #     text_frame2 = stitle.text_frame

    p2 = text_frame2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    p2.text = f'Generated: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}'
    p2.font.bold = False
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(0x30, 0x9c, 0x4f)
    p2.font.name = 'Arial'

    p3 = text_frame2.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.text = f'Software Version: {VERSION}'
    p3.font.bold = False
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(0x30, 0x9c, 0x4f)
    p3.font.name = 'Arial'

    return prs 


def subtitle_header(slide, title: str):
    """ Creates subtitle under main slide title """
    top = Inches(0.61)
    left = Inches(0.42)
    width = height = Inches(0.5)
    txtbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = txtbox.text_frame 

    p = text_frame.paragraphs[0]
    p.text = title 
    p.font.bold = False 
    p.font.size = Pt(22)
    p.font.name = 'Times New Roman'

    return slide


def make_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
    slide = fund_title_header(slide, 'Explanation of Analysis', include_time=False)

    if os.path.exists('metric_explanation.txt'):
        filer = open('metric_explanation.txt', 'r')
        content = filer.readlines()
        content2 = []
        for cont in content:
            c = cont.split('\r\n')[0]
            content2.append(c)
        content = content2
        filer.close()

        top = Inches(0.81)
        left = Inches(0.42)
        width = Inches(9)
        height = Inches(6)
        txtbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = txtbox.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = content[0] 
        p.font.size = Pt(12)
        p.font.bold = True
        for i in range(1,len(content)):
            p = text_frame.add_paragraph()
            p.text = content[i]
            if i == 3:
                p.font.size = Pt(12)
                p.font.bold = True
            else:
                p.font.size = Pt(10)
                p.font.bold = False

    else:
        print("WARNING - file 'metric_explanation.txt' not found.")

    return prs


def make_MCI_slides(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
    slide = fund_title_header(slide, 'Market Composite Index')
    
    content = f'output/temp/MCI.png'
    if os.path.exists(content):
        left = Inches(3)
        top = Inches(1.27)
        height = Inches(6)
        width = Inches(7.3)
        slide.shapes.add_picture(content, left, top, height=height, width=width)

    return prs


def make_fund_slides(prs, analysis: dict):
    funds = analysis.keys()
    for fund in funds:
        prs = add_fund_content(prs, fund, analysis)

    return prs


def add_fund_content(prs, fund: str, analysis: dict):
    # slide = prs.slides.add_slide(prs.slide_layouts[PRES_TITLE_SLIDE])
    # title = slide.shapes.title
    # text_frame = title.text_frame
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
    top = Inches(2.5)
    left = Inches(4)
    width = Inches(5)
    height = Inches(2)
    txtbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = txtbox.text_frame

    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = f'{fund}'
    p.font.bold = True
    p.font.size = Pt(54)
    p.font.name = 'Arial'

    p2 = text_frame.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = f"Dates Covered: {analysis[fund]['dates_covered']['start']}  :  {analysis[fund]['dates_covered']['end']}"
    p2.font.bold = False
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0x74, 0x3c, 0xe6)
    p2.font.name = 'Arial'

    # stitle = slide.placeholders[1]
    # text_frame2 = stitle.text_frame
    # p3 = text_frame2.paragraphs[0]
    # p3.text = ' '

    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
    indexes = []
    indexes.append(len(prs.slides) - 1)

    slide = fund_title_header(slide, fund)
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
    slide = fund_title_header(slide, fund)
    indexes.append(len(prs.slides)-1)

    content_dir = f'output/temp/{fund}/'
    if os.path.exists(content_dir):
        content = content_dir + '*.png'
        pics = glob.glob(content)
        prs = format_plots(prs, indexes, pics)

    return prs


def fund_title_header(slide, fund: str, include_time=True):
    left = Inches(0) #Inches(3.86)
    top = Inches(0)
    width = height = Inches(0.5)
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    #p = tf.add_paragraph()
    p = tf.paragraphs[0]
    p.text = fund 
    p.font.size = Pt(36)
    p.font.name = 'Arial'
    p.font.bold = True

    if include_time:
        p = tf.add_paragraph()
        p.font.size = Pt(14)
        p.font.bold = False
        p.font.name = 'Arial'
        p.text = str(datetime.now().strftime("%Y-%m-%d %H:%M:%S"))

    return slide


def format_plots(prs, slide_indices: list, globs: list):
    parts = windows_compatible_file_parse(globs[0])

    header = parts[0] + '/' + parts[1] + '/' + parts[2] + '/'

    for globber in globs:

        globbed = windows_compatible_file_parse(globber)
        part = globbed[3]

        if 'cluster' in part:
            left = Inches(0)
            top = Inches(1.1)
            height = Inches(3.0)
            width = Inches(4.5)
            prs.slides[slide_indices[0]].shapes.add_picture(header+part, left, top, height=height, width=width)

        if 'macd_bar' in part:
            left = Inches(4.5)
            top = Inches(1.1)
            height = Inches(3.0)
            width = Inches(4.5)
            prs.slides[slide_indices[0]].shapes.add_picture(header+part, left, top, height=height, width=width)

        if 'simple_moving_averages' in part:
            left = Inches(0.0)
            top = Inches(4.1)
            height = Inches(3.0)
            width = Inches(4.5)
            prs.slides[slide_indices[0]].shapes.add_picture(header+part, left, top, height=height, width=width)

        if 'obv' in part:
            left = Inches(4.5)
            top = Inches(4.1)
            height = Inches(3.0)
            width = Inches(4.5)
            prs.slides[slide_indices[0]].shapes.add_picture(header+part, left, top, height=height, width=width)

        ### Slide #2

        if 'relative_strength' in part:
            left = Inches(0)
            top = Inches(1.1)
            height = Inches(3.0)
            width = Inches(4.5)
            prs.slides[slide_indices[1]].shapes.add_picture(header+part, left, top, height=height, width=width)

        if 'exp_moving_averages' in part:
            left = Inches(4.5)
            top = Inches(1.1)
            height = Inches(3.0)
            width = Inches(4.5)
            prs.slides[slide_indices[1]].shapes.add_picture(header+part, left, top, height=height, width=width)

        if 'swing_trades' in part:
            left = Inches(0.0)
            top = Inches(4.1)
            height = Inches(3.0)
            width = Inches(4.5)
            prs.slides[slide_indices[1]].shapes.add_picture(header+part, left, top, height=height, width=width)

        if 'head_and_shoulders' in part:
            left = Inches(4.5)
            top = Inches(4.1)
            height = Inches(3.0)
            width = Inches(4.5)
            prs.slides[slide_indices[1]].shapes.add_picture(header+part, left, top, height=height, width=width)

    return prs 



def slide_creator(year: str, analysis: dict, version: str):
    """ High-level function for converting inventors spreadsheet to slides """

    print("Starting presentation creation.")

    prs = title_presentation(year, VERSION=version)
    prs = make_intro_slide(prs)
    prs = make_MCI_slides(prs)
    prs = make_fund_slides(prs, analysis)

    if not os.path.exists('output/'):
        os.mkdir('output/')
        
    prs.save(f'output/Financial_Analysis_{year}.pptx')
    print("Presentation created.")
